#!/usr/bin/env python3
import argparse
import os
import sys
import zipfile
from xml.etree import ElementTree as ET

# --- Gracefully Import Libraries & Provide Feedback ---
print("🐍 Checking for required Python libraries...")
try:
    import pypandoc
    HAS_PANDOC = True
    print("  > Found 'pypandoc' for high-quality conversion.")
except ImportError:
    HAS_PANDOC = False

try:
    from docx import Document
    HAS_DOCX = True
    print("  > Found 'python-docx' for .docx creation.")
except ImportError:
    HAS_DOCX = False

try:
    from mdutils.mdutils import MdUtils
    HAS_MDUTILS = True
    print("  > Found 'mdutils' for .md creation.")
except ImportError:
    HAS_MDUTILS = False

try:
    from PyRTF3.document import RTF, Section, Paragraph
    from PyRTF3.styles import TextStyle, ParagraphStyle
    HAS_RTF3 = True
    print("  > Found 'PyRTF3' for .rtf creation.")
except ImportError:
    HAS_RTF3 = False
    
try:
    from pptx import Presentation
    HAS_PPTX = True
    print("  > Found 'python-pptx' for core extraction.")
except ImportError:
    HAS_PPTX = False

print("-" * 20)


# --- XML Namespace Map ---
ns = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def extract_text_from_pptx(pptx_file, include_comments=True):
    """
    Extracts text and comments from a PowerPoint file, sorting text
    from shapes top-to-bottom and preserving internal line breaks.
    """
    prs = Presentation(pptx_file)
    extracted_data = []
    
    all_comments = {}
    if include_comments:
        try:
            with zipfile.ZipFile(pptx_file, 'r') as zf:
                for i in range(1, len(prs.slides) + 1):
                    rel_file_path = f'ppt/slides/_rels/slide{i}.xml.rels'
                    if rel_file_path in zf.namelist():
                        with zf.open(rel_file_path) as rel_file:
                            rel_tree = ET.parse(rel_file)
                            for rel in rel_tree.findall(".//r:Relationship", ns):
                                if "comments" in rel.get('Target'):
                                    comment_file_path = 'ppt/comments/' + os.path.basename(rel.get('Target'))
                                    if comment_file_path in zf.namelist():
                                        slide_comments = all_comments.setdefault(i, [])
                                        with zf.open(comment_file_path) as cf:
                                            comment_tree = ET.parse(cf)
                                            for comment in comment_tree.findall('.//p:cm', ns):
                                                text = comment.find('.//p:text', ns).text or ""
                                                author = comment.find('.//p:authorLst/p:author', ns).get('name', 'Unknown Author')
                                                slide_comments.append(f"{author}: {text}")
        except Exception as e:
            print(f"Warning: Could not extract comments. Error: {e}", file=sys.stderr)

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        shapes = sorted([s for s in slide.shapes if s.has_text_frame], key=lambda s: s.top)
        
        slide_text_parts = []
        for shape in shapes:
            shape_text = []
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text:
                    shape_text.append(paragraph.text)
            
            slide_text_parts.append("  \n".join(shape_text))

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                slide_text_parts.append("\n--- Speaker Notes ---\n" + notes_text)
        
        slide_text = "\n\n".join(filter(None, slide_text_parts)).strip()
        comments = all_comments.get(slide_number, [])
        extracted_data.append((slide_number, slide_text, comments))

    return extracted_data

def generate_markdown_string_for_pandoc(data):
    """Generates a raw Markdown string, intended for Pandoc conversion."""
    md_content = ""
    for slide_number, slide_text, comments in data:
        md_content += f"## Slide {slide_number}\n\n"
        md_content += f"{slide_text}\n\n"
        if comments:
            md_content += "### Comments\n\n"
            for comment in comments:
                md_content += f"* {comment}\n"
            md_content += "\n"
        md_content += "---\n\n"
    return md_content

def save_with_pandoc(md_string, output_file, output_format):
    """Saves content using Pandoc for high-quality conversion."""
    pypandoc.convert_text(md_string, output_format, format='md', outputfile=output_file)
    print(f"✅ Successfully converted to '{output_file}' using Pandoc.")

def save_as_markdown_native(data, output_file, wrap_width):
    """Saves content as a native Markdown file using mdutils, with optional text wrapping."""
    md_file = MdUtils(file_name=output_file.replace('.md', ''), title="Presentation Content")
    for slide_number, slide_text, comments in data:
        md_file.new_header(level=1, title=f"Slide {slide_number}")
        # Apply wrap_width to the paragraph containing the main slide text
        md_file.new_paragraph(slide_text, wrap_width=wrap_width)
        if comments:
            md_file.new_header(level=2, title="Comments")
            md_file.new_list(comments)
        
        md_file.new_line('---')

    md_file.create_md_file()
    print(f"✅ Successfully saved to '{output_file}' using mdutils.")

def save_as_docx_native(data, output_file):
    """Fallback to save as .docx using python-docx."""
    doc = Document()
    doc.add_heading("Presentation Content", level=0)
    for slide_number, slide_text, comments in data:
        doc.add_heading(f"Slide {slide_number}", level=1)
        doc.add_paragraph(slide_text)
        if comments:
            doc.add_heading("Comments", level=2)
            for comment in comments:
                doc.add_paragraph(comment, style='List Bullet')
    doc.save(output_file)
    print(f"✅ Successfully saved to '{output_file}' using python-docx (fallback).")

def save_as_rtf_native(data, output_file):
    """Fallback to save as .rtf using PyRTF3."""
    doc = RTF()
    section = Section()
    doc.Sections.append(section)
    h1 = ParagraphStyle(TextStyle(bold=True, size=32), spacing_after=240)
    h2 = ParagraphStyle(TextStyle(bold=True, size=24), spacing_before=120, spacing_after=120)

    for slide_number, slide_text, comments in data:
        section.append(Paragraph(f'Slide {slide_number}', h1))
        section.append(Paragraph(slide_text))
        if comments:
            section.append(Paragraph('Comments', h2))
            for comment in comments:
                section.append(Paragraph(f'- {comment}'))
        section.append(Paragraph('---', ParagraphStyle(spacing_before=480, spacing_after=480)))
        
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(doc.ToString())
    print(f"✅ Successfully saved to '{output_file}' using PyRTF3 (fallback).")


def main():
    """Main function to parse arguments and run the extraction with fallbacks."""
    if not HAS_PPTX:
        print("❌ FATAL ERROR: The core library 'python-pptx' is not installed.", file=sys.stderr)
        print("Please run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    parser = argparse.ArgumentParser(description="Robustly extract text and comments from a .pptx file.", formatter_class=argparse.RawTextHelpFormatter)
    
    parser.add_argument("input_file", help="Path to the input .pptx file.")
    parser.add_argument("output_format", nargs='?', default='docx', choices=['docx', 'md', 'rtf'], 
                        help="The desired output format (default: docx).")
    
    parser.add_argument("--no-comments", action="store_true", help="Exclude comments from the extraction.")
    
    parser.add_argument("--output-lib", default='auto', choices=['auto', 'pandoc', 'native'],
                        help="""Specify the output library to use:
  auto   - (Default) Use pandoc if available, otherwise fallback to native.
  pandoc - Force use of pandoc.
  native - Force use of native python libraries (python-docx, etc.).""")
    
    parser.add_argument("--wrap-text", type=int, default=0, metavar='WIDTH',
                        help="""Wrap text in markdown output at a given character width.
Default is 0 (no wrapping).""")

    args = parser.parse_args()

    if not os.path.exists(args.input_file):
        print(f"❌ ERROR: Input file '{args.input_file}' not found.", file=sys.stderr)
        sys.exit(1)

    base_name = os.path.splitext(args.input_file)[0]
    output_file = f"{base_name}.{args.output_format}"

    print(f"🔬 Extracting content from '{args.input_file}'...")
    extracted_data = extract_text_from_pptx(args.input_file, include_comments=not args.no_comments)
    print("  > Extraction complete.")

    print(f"💾 Saving to '{output_file}'...")
    
    use_pandoc = (args.output_lib == 'auto' and HAS_PANDOC) or (args.output_lib == 'pandoc')

    if use_pandoc and args.output_format != 'md':
        if not HAS_PANDOC:
            print("❌ ERROR: --output-lib set to 'pandoc', but 'pypandoc' is not installed.", file=sys.stderr)
            sys.exit(1)
        
        markdown_content = generate_markdown_string_for_pandoc(extracted_data)
        save_with_pandoc(markdown_content, output_file, args.output_format)
    
    else: # Fallback to native libraries
        if args.output_format == 'md':
            if HAS_MDUTILS:
                save_as_markdown_native(extracted_data, output_file, wrap_width=args.wrap_text)
            else:
                print("❌ ERROR: To save as .md, the required library is missing.", file=sys.stderr)
                print("Please run: pip install mdutils", file=sys.stderr)
                sys.exit(1)
        
        elif args.output_format == 'docx':
            if HAS_DOCX:
                save_as_docx_native(extracted_data, output_file)
            else:
                print("❌ ERROR: To save as .docx, the required library is missing.", file=sys.stderr)
                print("Please run: pip install python-docx", file=sys.stderr)
                sys.exit(1)

        elif args.output_format == 'rtf':
            if HAS_RTF3:
                save_as_rtf_native(extracted_data, output_file)
            else:
                print("❌ ERROR: To save as .rtf, the required library is missing.", file=sys.stderr)
                print("Please run: pip install PyRTF3", file=sys.stderr)
                sys.exit(1)

if __name__ == "__main__":
    main()
